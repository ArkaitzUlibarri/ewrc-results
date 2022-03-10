import os
import json
from config import app
from pptx import Presentation
from pptx.util import Cm
from pptx.util import Pt
from database import driver_stats
from database import team_stats


def create_table(body_data):
    rows, cols = len(body_data) + 1, len(body_data[0])
    left, top, width, height = Cm(1), Cm(4), Cm(rows * 24 / 14), Cm(cols * 11 / 6)
    return shapes.add_table(rows, cols, left, top, width, height).table


def write_table(ppt_table, headings, body_data):
    for heading_index, heading_item in enumerate(headings, start=0):
        cell = ppt_table.cell(0, heading_index)
        cell.text = heading_item
        set_cell_font_size(cell, 14)

    for body_index, body_item in enumerate(body_data, start=1):
        for col_index, text in enumerate(body_item, start=0):
            cell = ppt_table.cell(body_index, col_index)
            cell.text = str(text)
            set_cell_font_size(cell, 12)


def set_cell_font_size(cell, size):
    for paragraph in cell.text_frame.paragraphs:
        for paragraph_run in paragraph.runs:
            paragraph_run.font.size = Pt(size)


def get_points(position, points_dict):
    points = '0'

    for item in points_dict:
        if position in item['position']:
            points = item['points']

    return points


def get_driver_season_stats(driver_id, results, points_dict):
    output = dict()
    output['driver_id'] = driver_id
    output['results'] = list()

    starts = 0
    wins = 0
    podiums = 0
    top5 = 0
    dnfs = 0
    total_points = 0
    for item_index, item in enumerate(results, start=1):
        starts += 1
        position = item['result']
        season_event_id = item['ID']
        points = get_points(position, points_dict)

        if not position.isnumeric():
            dnfs += 1
        else:
            if int(position) == 1:
                wins += 1
            if int(position) <= 3:
                podiums += 1
            if int(position) <= 5:
                top5 += 1

        total_points += int(points)

        output['results'].append({"season_event_id": season_event_id, "position": position, "points": points})
    output['starts'] = starts
    output['wins'] = wins
    output['podiums'] = podiums
    output['top5'] = top5
    output['dnfs'] = dnfs
    output['total_points'] = total_points

    return output


# Clear console
os.system("cls")

package_dir = os.path.abspath(os.path.dirname(__file__))
db_path = os.path.join(package_dir, 'database', app.database + '.db')

# Config
season = str(1997)

# Queries - Drivers
drivers_scratchs_stats = driver_stats.drivers_scratchs(db_path, season)
drivers_leaders_stats = driver_stats.drivers_leaders(db_path, season)
drivers_winners_stats = driver_stats.drivers_winners(db_path, season)
drivers_podiums_stats = driver_stats.drivers_podiums(db_path, season)
full_season_winners = driver_stats.season_winners(db_path, season)

# Queries - Teams
teams_scratchs_stats = team_stats.teams_scratchs(db_path, season)
teams_leaders_stats = team_stats.teams_leaders(db_path, season)
teams_winners_stats = team_stats.teams_winners(db_path, season)
teams_podiums_stats = team_stats.teams_podiums(db_path, season)

# Queries - Points System
driver_points_system_dict = json.loads(driver_stats.championship_points_system(db_path, season, 'drivers'))

# Obtain all the drivers which have scored points
lowest_position = len(driver_points_system_dict)
drivers_in_points_list = driver_stats.drivers_in_points(db_path, season, lowest_position)

full_results = list()
for index, row in enumerate(drivers_in_points_list, start=1):
    results_by_driver_list = driver_stats.full_results_by_driver(db_path, season, row['driver_id'])
    driver_results_dict = get_driver_season_stats(row['driver_id'], results_by_driver_list, driver_points_system_dict)
    full_results.append(driver_results_dict)
full_results = sorted(full_results, key=lambda k: k['total_points'], reverse=True)
# print(full_results)

# TODO: Drivers on the podium + official teams - SLIDE ON PPT
# TODO: Championship Table Slide - Only TOP 10

# create PPT
prs = Presentation()
title_slide_layout = prs.slide_layouts[5]

# 1 Slide - Season Winners
if len(full_season_winners):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    text_frame = shapes[0].text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'WORLD RALLY CHAMPIONSHIP ' + season

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(40)

    table = create_table(full_season_winners)
    write_table(table, ('#', 'Edition', 'Rally', 'Winner', 'Car', 'Team'), full_season_winners)

# DRIVER

# 2 Slide - Win stats
if len(drivers_winners_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'DRIVER STATS ' + season

    table = create_table(drivers_winners_stats)
    write_table(table, ('Driver', 'Wins'), drivers_winners_stats)

# 3 Slide - Podium stats
if len(drivers_podiums_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'DRIVER STATS ' + season

    table = create_table(drivers_podiums_stats)
    write_table(table, ('Driver', 'Podiums'), drivers_podiums_stats)

# 4 Slide - Scratchs stats
if len(drivers_scratchs_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'DRIVER STATS ' + season

    table = create_table(drivers_scratchs_stats)
    write_table(table, ('Driver', 'Scratchs'), drivers_scratchs_stats)

# 5 Slide - Leaders stats
if len(drivers_leaders_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'DRIVER STATS ' + season

    table = create_table(drivers_leaders_stats)
    write_table(table, ('Driver', 'Leaders'), drivers_leaders_stats)

# TEAM

# 6 Slide - Win stats
if len(teams_winners_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'TEAM STATS ' + season

    table = create_table(teams_winners_stats)
    write_table(table, ('Car', 'Team', 'Wins'), teams_winners_stats)

# 7 Slide - Podium stats
if len(teams_podiums_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'TEAM STATS ' + season

    table = create_table(teams_podiums_stats)
    write_table(table, ('Car', 'Team', 'Podiums'), teams_podiums_stats)

# 8 Slide - Scratchs stats
if len(teams_scratchs_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'TEAM STATS ' + season

    table = create_table(teams_scratchs_stats)
    write_table(table, ('Team', 'Scratchs'), teams_scratchs_stats)

# 9 Slide - Leaders stats
if len(teams_leaders_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'TEAM STATS ' + season

    table = create_table(teams_leaders_stats)
    write_table(table, ('Team', 'Leaders'), teams_leaders_stats)

# Save PPT
package_dir = os.path.abspath(os.path.dirname(__file__))
export_folder = os.path.join(package_dir, 'storage', 'exports')
if not os.path.exists(export_folder):
    os.makedirs(export_folder)

export_path = os.path.join(package_dir, 'storage', 'exports', 'WRC_' + season + '.pptx')
prs.save(export_path)

print('Finished ' + export_path + ' export')

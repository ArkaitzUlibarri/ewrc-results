import json
import logging
import os

from pptx import Presentation
from pptx.util import Pt

import definitions
from services.ppt import driver_service
from services import event_service
from services import point_service
from services.ppt import ppt_service
from services.ppt import team_service


def get_points(position, points_dict):
    points = '0'

    for item in points_dict:
        if position in item['position']:
            points = item['points']

    return points


def get_driver_standings(driver_data, results, points_dict):
    output = dict()
    output['driver_id'] = driver_data['driver_id']
    output['driver'] = driver_data['fullname']
    output['results'] = list()
    output['starts'] = 0
    output['wins'] = 0
    output['podiums'] = 0
    output['top5'] = 0
    output['dnfs'] = 0
    output['total_points'] = 0
    for entry_index, entry in enumerate(results, start=1):
        if entry['result'] is None:
            position = '-'
            points = '0'
        else:
            output['starts'] += 1
            if entry['result'].isnumeric():
                position = entry['result']
                points = get_points(position, points_dict)
                if int(position) == 1:
                    output['wins'] += 1
                if int(position) <= 3:
                    output['podiums'] += 1
                if int(position) <= 5:
                    output['top5'] += 1
            else:
                position = 'R'
                points = '0'
                output['dnfs'] += 1
        output['total_points'] += int(points)
        output['results'].append({
            "event_id": entry['event_id'],
            "season_order": entry['season_order'],
            "result": entry['result'],
            "position": position,
            "points": points
        })

    return output


def get_team_standings(team_data, results, points_dict):
    output = dict()
    output['car'] = team_data['car']
    output['team'] = team_data['team']
    output['results'] = dict()
    output['starts'] = 0
    output['wins'] = 0
    output['podiums'] = 0
    output['top5'] = 0
    output['dnfs'] = 0
    output['total_points'] = 0
    for entry_index, entry in enumerate(results, start=1):
        output['starts'] += 1
        if entry['result'].isnumeric():
            position = entry['result']
            points = get_points(position, points_dict)
            if int(position) == 1:
                output['wins'] += 1
            if int(position) <= 3:
                output['podiums'] += 1
            if int(position) <= 5:
                output['top5'] += 1
        else:
            points = '0'
            output['dnfs'] += 1
        output['total_points'] += int(points)
        if entry['event_id'] in output['results'].keys():
            output['results'][entry['event_id']] += int(points)
        else:
            output['results'][entry['event_id']] = int(points)

    return output


# Clear console
os.system("cls")

# Config
season = str(1998)

full_season_winners = driver_service.get_full_season_winners(season)
season_events = event_service.get_season_events(season)

# DRIVERS

# Queries - Drivers
drivers_scratchs_stats = driver_service.get_season_scratchs(season)
drivers_leaders_stats = driver_service.get_season_leaders(season)
drivers_winners_stats = driver_service.get_season_winners(season)
drivers_podiums_stats = driver_service.get_season_podiums(season)

# Queries - Points System
drivers_points_system_dict = json.loads(point_service.championship_points_system(season, 'drivers'))

# Obtain all the drivers which have scored points
lowest_position = len(drivers_points_system_dict)
drivers_in_points_list = driver_service.get_drivers_in_points(season, lowest_position)

# Full championship standings
full_championship_standings = list()
for index, driver_data in enumerate(drivers_in_points_list, start=1):
    results_by_driver_list = driver_service.get_driver_season_results(season, driver_data['driver_id'])
    driver_results_dict = get_driver_standings(driver_data, results_by_driver_list, drivers_points_system_dict)
    full_championship_standings.append(driver_results_dict)
full_championship_standings = sorted(full_championship_standings, key=lambda k: k['total_points'], reverse=True)

limit = 10
driver_standings = []
driver_standings_header = ('#', 'Drivers',)
for index, event in enumerate(season_events, start=1):
    driver_standings_header += (str(index),)
driver_standings_header += ('Points',)

for index, driver_result in enumerate(full_championship_standings, start=1):
    row = (index, driver_result['driver'],)
    for result in driver_result['results']:
        row += (result['points'] + ' (' + result['position'] + ')',)
    row += (driver_result['total_points'],)
    driver_standings.append(row)
    if index == limit:
        break

# TEAMS

# Queries - Teams
teams_scratchs_stats = team_service.get_season_scratchs(season)
teams_leaders_stats = team_service.get_season_leaders(season)
teams_winners_stats = team_service.get_season_winners(season)
teams_podiums_stats = team_service.get_season_podiums(season)

# Queries - Points System
teams_points_system_dict = json.loads(point_service.championship_points_system(season, 'manufacturers'))

# Obtain all the drivers which have scored points
lowest_position = len(teams_points_system_dict)
teams_in_points_list = team_service.get_teams_in_points(season, lowest_position)

# FIXME: Team standings
# Full championship standings
full_championship_standings = list()
for index, team_data in enumerate(teams_in_points_list, start=1):
    results_by_team_list = team_service.get_team_season_results(season, team_data['team'])
    driver_results_dict = get_team_standings(team_data, results_by_team_list, teams_points_system_dict)
    full_championship_standings.append(driver_results_dict)
full_championship_standings = sorted(full_championship_standings, key=lambda k: k['total_points'], reverse=True)

team_standings = []
team_standings_header = ('#', 'Teams',)
for index, event in enumerate(season_events, start=1):
    team_standings_header += (str(index),)
team_standings_header += ('Points',)

for index, team_result in enumerate(full_championship_standings, start=1):
    row = (index, team_result['team'],) + tuple(team_result['results'].values()) + (team_result['total_points'],)
    team_standings.append(row)

# TODO: Stats foreach TOP 10 driver
# TODO: Stats per team

# create PPT
prs = Presentation()
title_slide_layout = prs.slide_layouts[5]

# 1 Slide - Season Winners
if len(full_season_winners):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    text_frame = shapes[0].text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'WORLD RALLY CHAMPIONSHIP ' + season

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(40)

    table = ppt_service.create_table(shapes, full_season_winners)
    ppt_service.write_table(table, ('#', 'Edition', 'Rally', 'Winner', 'Car', 'Team'), full_season_winners)

# DRIVER

# 2 Slide - Win stats
if len(drivers_winners_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'DRIVER STATS ' + season

    table = ppt_service.create_table(shapes, drivers_winners_stats)
    ppt_service.write_table(table, ('Driver', 'Wins'), drivers_winners_stats)

# 3 Slide - Podium stats
if len(drivers_podiums_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'DRIVER STATS ' + season

    table = ppt_service.create_table(shapes, drivers_podiums_stats)
    ppt_service.write_table(table, ('Driver', 'Podiums'), drivers_podiums_stats)

# 4 Slide - Scratchs stats
if len(drivers_scratchs_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'DRIVER STATS ' + season

    table = ppt_service.create_table(shapes, drivers_scratchs_stats)
    ppt_service.write_table(table, ('Driver', 'Scratchs'), drivers_scratchs_stats)

# 5 Slide - Leaders stats
if len(drivers_leaders_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'DRIVER STATS ' + season

    table = ppt_service.create_table(shapes, drivers_leaders_stats)
    ppt_service.write_table(table, ('Driver', 'Leaders'), drivers_leaders_stats)

# TEAM

# 6 Slide - Win stats
if len(teams_winners_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'TEAM STATS ' + season

    table = ppt_service.create_table(shapes, teams_winners_stats)
    ppt_service.write_table(table, ('Car', 'Team', 'Wins'), teams_winners_stats)

# 7 Slide - Podium stats
if len(teams_podiums_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'TEAM STATS ' + season

    table = ppt_service.create_table(shapes, teams_podiums_stats)
    ppt_service.write_table(table, ('Car', 'Team', 'Podiums'), teams_podiums_stats)

# 8 Slide - Scratchs stats
if len(teams_scratchs_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'TEAM STATS ' + season

    table = ppt_service.create_table(shapes, teams_scratchs_stats)
    ppt_service.write_table(table, ('Team', 'Scratchs'), teams_scratchs_stats)

# 9 Slide - Leaders stats
if len(teams_leaders_stats):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'TEAM STATS ' + season

    table = ppt_service.create_table(shapes, teams_leaders_stats)
    ppt_service.write_table(table, ('Team', 'Leaders'), teams_leaders_stats)

# 10 Driver Championship Standings (Top 10)
if len(driver_standings):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'DRIVER CHAMPIONSHIP STANDINGS ' + season

    table = ppt_service.create_table(shapes, driver_standings)
    ppt_service.write_table(table, driver_standings_header, driver_standings)

# 11 Teams Championship Standings
if len(team_standings):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'TEAM CHAMPIONSHIP STANDINGS ' + season

    table = ppt_service.create_table(shapes, team_standings)
    ppt_service.write_table(table, team_standings_header, team_standings)

# Save PPT
if not os.path.exists(definitions.EXPORT_FOLDER):
    os.makedirs(definitions.EXPORT_FOLDER)

export_path = os.path.join(definitions.ROOT_DIR, '../storage', 'exports', 'WRC_' + season + '.pptx')
prs.save(export_path)

logging.info('Finished ' + export_path + ' export')

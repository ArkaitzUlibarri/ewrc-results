import os
import json
from config import app
from pptx import Presentation
from pptx.util import Cm
from database.driver_stats import *
from database.team_stats import *


def create_table(data):
    rows, cols = len(data) + 1, len(data[0])
    left, top, width, height = Cm(1), Cm(4), Cm(rows * 24 / 14), Cm(cols * 11 / 6)
    return shapes.add_table(rows, cols, left, top, width, height).table


def write_table_headings(ppt_table, headings):
    for index, heading in enumerate(headings, start=0):
        ppt_table.cell(0, index).text = heading


def write_table_body(ppt_table, data):
    for row_index, row in enumerate(data, start=1):
        for col_index, cell in enumerate(row, start=0):
            ppt_table.cell(row_index, col_index).text = str(cell)


def get_points(result, points_dict):
    points = '0'

    for item in points_dict:
        if result in item['position']:
            points = item['points']

    return points


def get_driver_season_stats(results, points_dict):
    starts = 0
    wins = 0
    podiums = 0
    top5 = 0
    dnfs = 0
    total_points = 0
    for row_index, row in enumerate(results, start=1):
        starts += 1
        result = row['result']
        points = get_points(result, points_dict)

        if not result.isnumeric():
            dnfs += 1
        else:
            if int(result) == 1:
                wins += 1
            if int(result) <= 3:
                podiums += 1
            if int(result) <= 5:
                top5 += 1
        
        total_points += int(points)
        print("Result: " + result + ", Points: " + points)
    print("Starts: " + str(starts)
          + ", Wins: " + str(wins)
          + ", Podiums: " + str(podiums)
          + ", TOP5: " + str(top5)
          + ", DNFs: " + str(dnfs)
          + ", Points: " + str(total_points))


# Clear console
os.system("cls")

package_dir = os.path.abspath(os.path.dirname(__file__))
db_path = os.path.join(package_dir, 'database', app.database + '.db')

# Config
season = str(1997)

# Queries - Drivers
drivers_scratchs_stats = drivers_scratchs(db_path, season)
drivers_leaders_stats = drivers_leaders(db_path, season)
drivers_winners_stats = drivers_winners(db_path, season)
drivers_podiums_stats = drivers_podiums(db_path, season)
full_season_winners = season_winners(db_path, season)

# Queries - Teams
teams_winners_stats = teams_winners(db_path, season)
teams_podiums_stats = teams_podiums(db_path, season)

# Queries - Points System
driver_points_system_dict = json.loads(championship_points_system(db_path, season, 'drivers'))

# TODO: Obtain all the drivers which have scored points
lowest_position = len(driver_points_system_dict)

# Driver Points
# TODO: Drivers on the podium + official teams
full_results_by_driver = full_results_by_driver(db_path, season, 848)
get_driver_season_stats(full_results_by_driver, driver_points_system_dict)

# create PPT
prs = Presentation()
layout = prs.slide_layouts[5]

# First Slide - Season Winners
slide = prs.slides.add_slide(layout)
shapes = slide.shapes
shapes.title.text = 'WORLD RALLY CHAMPIONSHIP ' + season

table = create_table(full_season_winners)
write_table_headings(table, ('#', 'Edition', 'Rally', 'Winner', 'Car', 'Team'))
write_table_body(table, full_season_winners)

# Second slide - Win stats
slide = prs.slides.add_slide(layout)
shapes = slide.shapes
shapes.title.text = 'STATS ' + season

table = create_table(drivers_winners_stats)
write_table_headings(table, ('Driver', 'Wins'))
write_table_body(table, drivers_winners_stats)

# Third slide - Podium stats
slide = prs.slides.add_slide(layout)
shapes = slide.shapes
shapes.title.text = 'STATS ' + season

table = create_table(drivers_podiums_stats)
write_table_headings(table, ('Driver', 'Podiums'))
write_table_body(table, drivers_podiums_stats)

# Fourth slide - Scratchs stats
slide = prs.slides.add_slide(layout)
shapes = slide.shapes
shapes.title.text = 'STATS ' + season

table = create_table(drivers_scratchs_stats)
write_table_headings(table, ('Driver', 'Scratchs'))
write_table_body(table, drivers_scratchs_stats)

# Fifth slide - Leaders stats
slide = prs.slides.add_slide(layout)
shapes = slide.shapes
shapes.title.text = 'STATS ' + season

table = create_table(drivers_leaders_stats)
write_table_headings(table, ('Driver', 'Leaders'))
write_table_body(table, drivers_leaders_stats)

# Sixth slide - Win stats
slide = prs.slides.add_slide(layout)
shapes = slide.shapes
shapes.title.text = 'STATS ' + season

table = create_table(teams_winners_stats)
write_table_headings(table, ('Car', 'Team', 'Wins'))
write_table_body(table, teams_winners_stats)

# Seven slide - Podium stats
slide = prs.slides.add_slide(layout)
shapes = slide.shapes
shapes.title.text = 'STATS ' + season

table = create_table(teams_podiums_stats)
write_table_headings(table, ('Car', 'Team', 'Podiums'))
write_table_body(table, teams_podiums_stats)

# Save PPT
package_dir = os.path.abspath(os.path.dirname(__file__))
export_folder = os.path.join(package_dir, 'storage', 'exports')
if not os.path.exists(export_folder):
    os.makedirs(export_folder)

export_path = os.path.join(package_dir, 'storage', 'exports', 'WRC_' + season + '.pptx')
prs.save(export_path)

print('Finished ' + export_path + ' export')
